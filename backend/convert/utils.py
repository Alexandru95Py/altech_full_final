import io
import os
import fitz  # PyMuPDF
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import zipfile


def convert_pdf_to_txt(uploaded_file):
    pdf_bytes = uploaded_file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")

    extracted_text = ""
    for page in doc:
        extracted_text += page.get_text()

    output_stream = io.BytesIO()
    output_stream.write(extracted_text.encode("utf-8"))
    output_stream.seek(0)
    return output_stream, "converted.txt", "text/plain"


def convert_pdf_to_docx(uploaded_file):
    input_stream = io.BytesIO(uploaded_file.read())
    input_stream.seek(0)

    temp_input_path = "temp_input.pdf"
    temp_output_path = "temp_output.docx"

    with open(temp_input_path, "wb") as f:
        f.write(input_stream.read())

    converter = Converter(temp_input_path)
    converter.convert(temp_output_path, start=0, end=None)
    converter.close()

    with open(temp_output_path, "rb") as out_f:
        output_stream = io.BytesIO(out_f.read())
    output_stream.seek(0)

    os.remove(temp_input_path)
    os.remove(temp_output_path)

    return output_stream, "converted.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def convert_pdf_to_pptx(uploaded_file):
    pdf_bytes = uploaded_file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("png")

        image_stream = io.BytesIO(img_bytes)
        image = Image.open(image_stream)

        png_stream = io.BytesIO()
        image.save(png_stream, format="PNG")
        png_stream.seek(0)

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(png_stream, Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)

    output_stream = io.BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)

    return output_stream, "converted.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def convert_pdf_to_images(uploaded_file, image_format="jpg"):
    import tempfile
    from PIL import Image
    import os

    print("🔍 Începem conversia PDF →", image_format.upper())

    pdf_bytes = uploaded_file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")

    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as zip_file:
        for page_num in range(len(doc)):
            try:
                print(f"📄 Pagina {page_num + 1}: generăm imagine")
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=150)

                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
                    temp_path = tmp_file.name
                    pix.save(temp_path)

                image = Image.open(temp_path)
                image.load()  # forțăm încărcarea completă în memorie

                if image_format == "jpg":
                    image = image.convert("RGB")

                final_img = io.BytesIO()
                image.save(final_img, format="JPEG" if image_format == "jpg" else "PNG")
                final_img.seek(0)

                filename = f"page_{page_num + 1}.{image_format}"
                zip_file.writestr(filename, final_img.read())
                print(f"✅ Adăugată în arhivă: {filename}")

                image.close()  # 👈 închidem imaginea
                os.remove(temp_path)  # 🔐 acum putem șterge

            except Exception as e:
                print(f"❌ Eroare la pagina {page_num + 1}: {e}")

    zip_buffer.seek(0)
    print("✅ Arhiva ZIP finalizată")
    return zip_buffer, f"converted_images_{image_format}.zip", "application/zip"

def convert_pdf_to_format(uploaded_file, target_format):
    if target_format == "txt":
        return convert_pdf_to_txt(uploaded_file)
    elif target_format == "docx":
        return convert_pdf_to_docx(uploaded_file)
    elif target_format == "pptx":
        return convert_pdf_to_pptx(uploaded_file)
    elif target_format == "jpg":
        return convert_pdf_to_images(uploaded_file, image_format="jpg")
    elif target_format == "png":
        return convert_pdf_to_images(uploaded_file, image_format="png")
    else:
        raise ValueError(f"Unsupported conversion format: {target_format}")